"""
文件生成器 - 支持多种文件格式的生成和处理
支持: 文档、表格、演示、图片、代码、数据、压缩等多种类型
"""
import os
import re
import json
import uuid
import base64
import zipfile
import tarfile
import gzip
from pathlib import Path
from datetime import datetime
from typing import Optional, Dict, Any, Tuple, List, Union
from enum import Enum
from dataclasses import dataclass

import pandas as pd

# ============================================================
# 可选依赖 - 按需导入
# ============================================================

# PDF 生成
try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False

# PPT 生成
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# Word 生成
try:
    from docx import Document
    from docx.shared import Inches as DocxInches, Pt as DocxPt
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

# 图片处理
try:
    from PIL import Image as PILImage
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

# YAML 处理
try:
    import yaml
    HAS_YAML = True
except ImportError:
    HAS_YAML = False

# Excel 处理 (openpyxl)
try:
    import openpyxl
    HAS_OPENPYXL = True
    print(f"[FileGenerator] openpyxl version: {openpyxl.__version__}")
except ImportError:
    HAS_OPENPYXL = False
    print("[FileGenerator] WARNING: openpyxl not installed, Excel generation may fail")


# ============================================================
# 文件类型定义
# ============================================================

class FileCategory(Enum):
    """文件类别"""
    DOCUMENT = "document"      # 文档
    SPREADSHEET = "spreadsheet"  # 表格
    PRESENTATION = "presentation"  # 演示
    IMAGE = "image"            # 图片
    AUDIO = "audio"            # 音频
    VIDEO = "video"            # 视频
    ARCHIVE = "archive"        # 压缩
    CODE = "code"              # 代码
    DATA = "data"              # 数据
    EXECUTABLE = "executable"  # 可执行
    FONT = "font"              # 字体
    DESIGN = "design"          # 设计
    OTHER = "other"            # 其他


@dataclass
class FileTypeInfo:
    """文件类型信息"""
    extension: str           # 扩展名
    category: FileCategory   # 类别
    mime_type: str          # MIME 类型
    keywords: List[str]     # 关键词（用于匹配）
    can_generate: bool      # 是否可生成
    description: str        # 描述


# 文件类型注册表
FILE_TYPES: Dict[str, FileTypeInfo] = {
    # ==================== 文档类 ====================
    "txt": FileTypeInfo("txt", FileCategory.DOCUMENT, "text/plain",
                        ["txt", "文本", "text", "纯文本"], True, "纯文本文件"),
    "docx": FileTypeInfo("docx", FileCategory.DOCUMENT, "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                         ["docx", "word", "文档", "文稿", "doc"], True, "Word 文档"),
    "pdf": FileTypeInfo("pdf", FileCategory.DOCUMENT, "application/pdf",
                        ["pdf"], True, "PDF 文档"),
    "rtf": FileTypeInfo("rtf", FileCategory.DOCUMENT, "application/rtf",
                        ["rtf", "富文本"], True, "富文本文档"),
    "odt": FileTypeInfo("odt", FileCategory.DOCUMENT, "application/vnd.oasis.opendocument.text",
                        ["odt", "opendocument"], False, "OpenDocument 文本"),
    "md": FileTypeInfo("md", FileCategory.DOCUMENT, "text/markdown",
                       ["md", "markdown"], True, "Markdown 文档"),

    # ==================== 表格类 ====================
    "xlsx": FileTypeInfo("xlsx", FileCategory.SPREADSHEET, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                         ["xlsx", "excel", "表格", "xls", "电子表格"], True, "Excel 表格"),
    "csv": FileTypeInfo("csv", FileCategory.SPREADSHEET, "text/csv",
                        ["csv"], True, "CSV 表格"),
    "ods": FileTypeInfo("ods", FileCategory.SPREADSHEET, "application/vnd.oasis.opendocument.spreadsheet",
                        ["ods"], False, "OpenDocument 表格"),

    # ==================== 演示类 ====================
    "pptx": FileTypeInfo("pptx", FileCategory.PRESENTATION, "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                         ["pptx", "ppt", "powerpoint", "演示", "幻灯片", "slides"], True, "PowerPoint 演示"),
    "odp": FileTypeInfo("odp", FileCategory.PRESENTATION, "application/vnd.oasis.opendocument.presentation",
                        ["odp"], False, "OpenDocument 演示"),

    # ==================== 图片类 - 位图 ====================
    "png": FileTypeInfo("png", FileCategory.IMAGE, "image/png",
                        ["png", "canvas", "画布", "绘图", "draw", "design", "设计图"], True, "PNG 图片"),
    "jpg": FileTypeInfo("jpg", FileCategory.IMAGE, "image/jpeg",
                        ["jpg", "jpeg", "图片", "image", "照片", "photo"], True, "JPEG 图片"),
    "gif": FileTypeInfo("gif", FileCategory.IMAGE, "image/gif",
                        ["gif", "动图"], True, "GIF 图片"),
    "bmp": FileTypeInfo("bmp", FileCategory.IMAGE, "image/bmp",
                        ["bmp", "bitmap"], True, "BMP 图片"),
    "webp": FileTypeInfo("webp", FileCategory.IMAGE, "image/webp",
                         ["webp"], True, "WebP 图片"),
    "tiff": FileTypeInfo("tiff", FileCategory.IMAGE, "image/tiff",
                         ["tiff", "tif"], True, "TIFF 图片"),
    "ico": FileTypeInfo("ico", FileCategory.IMAGE, "image/x-icon",
                        ["ico", "icon", "图标"], True, "图标文件"),

    # ==================== 图片类 - 矢量 ====================
    "svg": FileTypeInfo("svg", FileCategory.IMAGE, "image/svg+xml",
                        ["svg", "矢量"], True, "SVG 矢量图"),

    # ==================== 音频类 ====================
    "mp3": FileTypeInfo("mp3", FileCategory.AUDIO, "audio/mpeg",
                        ["mp3", "音频", "audio", "音乐"], False, "MP3 音频"),
    "wav": FileTypeInfo("wav", FileCategory.AUDIO, "audio/wav",
                        ["wav"], False, "WAV 音频"),
    "flac": FileTypeInfo("flac", FileCategory.AUDIO, "audio/flac",
                         ["flac"], False, "FLAC 音频"),
    "m4a": FileTypeInfo("m4a", FileCategory.AUDIO, "audio/mp4",
                        ["m4a", "aac"], False, "M4A 音频"),
    "ogg": FileTypeInfo("ogg", FileCategory.AUDIO, "audio/ogg",
                        ["ogg"], False, "OGG 音频"),

    # ==================== 视频类 ====================
    "mp4": FileTypeInfo("mp4", FileCategory.VIDEO, "video/mp4",
                        ["mp4", "视频", "video"], False, "MP4 视频"),
    "avi": FileTypeInfo("avi", FileCategory.VIDEO, "video/x-msvideo",
                        ["avi"], False, "AVI 视频"),
    "mkv": FileTypeInfo("mkv", FileCategory.VIDEO, "video/x-matroska",
                        ["mkv"], False, "MKV 视频"),
    "mov": FileTypeInfo("mov", FileCategory.VIDEO, "video/quicktime",
                        ["mov"], False, "MOV 视频"),
    "webm": FileTypeInfo("webm", FileCategory.VIDEO, "video/webm",
                         ["webm"], False, "WebM 视频"),

    # ==================== 压缩/归档 ====================
    "zip": FileTypeInfo("zip", FileCategory.ARCHIVE, "application/zip",
                        ["zip", "压缩", "archive", "打包"], True, "ZIP 压缩包"),
    "tar": FileTypeInfo("tar", FileCategory.ARCHIVE, "application/x-tar",
                        ["tar"], True, "TAR 归档"),
    "gz": FileTypeInfo("gz", FileCategory.ARCHIVE, "application/gzip",
                       ["gz", "gzip"], True, "GZIP 压缩"),
    "7z": FileTypeInfo("7z", FileCategory.ARCHIVE, "application/x-7z-compressed",
                       ["7z"], False, "7Z 压缩包"),
    "rar": FileTypeInfo("rar", FileCategory.ARCHIVE, "application/vnd.rar",
                        ["rar"], False, "RAR 压缩包"),

    # ==================== 代码/网页 ====================
    "html": FileTypeInfo("html", FileCategory.CODE, "text/html",
                         ["html", "htm", "网页", "前端", "页面", "web"], True, "HTML 网页"),
    "css": FileTypeInfo("css", FileCategory.CODE, "text/css",
                        ["css", "样式"], True, "CSS 样式表"),
    "js": FileTypeInfo("js", FileCategory.CODE, "application/javascript",
                       ["js", "javascript", "脚本"], True, "JavaScript 脚本"),
    "ts": FileTypeInfo("ts", FileCategory.CODE, "application/typescript",
                       ["ts", "typescript"], True, "TypeScript 脚本"),
    "py": FileTypeInfo("py", FileCategory.CODE, "text/x-python",
                       ["py", "python"], True, "Python 脚本"),
    "java": FileTypeInfo("java", FileCategory.CODE, "text/x-java-source",
                         ["java"], True, "Java 源码"),
    "c": FileTypeInfo("c", FileCategory.CODE, "text/x-c",
                      ["c"], True, "C 源码"),
    "cpp": FileTypeInfo("cpp", FileCategory.CODE, "text/x-c++",
                        ["cpp", "c++", "cxx"], True, "C++ 源码"),
    "go": FileTypeInfo("go", FileCategory.CODE, "text/x-go",
                       ["go", "golang"], True, "Go 源码"),
    "rs": FileTypeInfo("rs", FileCategory.CODE, "text/x-rust",
                       ["rs", "rust"], True, "Rust 源码"),
    "php": FileTypeInfo("php", FileCategory.CODE, "application/x-php",
                        ["php"], True, "PHP 脚本"),
    "rb": FileTypeInfo("rb", FileCategory.CODE, "text/x-ruby",
                       ["rb", "ruby"], True, "Ruby 脚本"),
    "sh": FileTypeInfo("sh", FileCategory.CODE, "application/x-sh",
                       ["sh", "bash", "shell"], True, "Shell 脚本"),
    "bat": FileTypeInfo("bat", FileCategory.CODE, "application/x-bat",
                        ["bat", "cmd", "batch"], True, "批处理脚本"),
    "ps1": FileTypeInfo("ps1", FileCategory.CODE, "application/x-powershell",
                        ["ps1", "powershell"], True, "PowerShell 脚本"),
    "vue": FileTypeInfo("vue", FileCategory.CODE, "text/x-vue",
                        ["vue"], True, "Vue 组件"),
    "jsx": FileTypeInfo("jsx", FileCategory.CODE, "text/jsx",
                        ["jsx", "react"], True, "JSX/React 组件"),
    "tsx": FileTypeInfo("tsx", FileCategory.CODE, "text/tsx",
                        ["tsx"], True, "TSX 组件"),

    # ==================== 数据格式 ====================
    "json": FileTypeInfo("json", FileCategory.DATA, "application/json",
                         ["json", "api"], True, "JSON 数据"),
    "xml": FileTypeInfo("xml", FileCategory.DATA, "application/xml",
                        ["xml"], True, "XML 数据"),
    "yaml": FileTypeInfo("yaml", FileCategory.DATA, "application/x-yaml",
                         ["yaml", "yml", "配置"], True, "YAML 数据"),
    "toml": FileTypeInfo("toml", FileCategory.DATA, "application/toml",
                         ["toml"], True, "TOML 数据"),
    "ini": FileTypeInfo("ini", FileCategory.DATA, "text/plain",
                        ["ini", "conf", "config"], True, "INI 配置"),
    "sql": FileTypeInfo("sql", FileCategory.DATA, "application/sql",
                        ["sql", "数据库", "database"], True, "SQL 脚本"),
    "log": FileTypeInfo("log", FileCategory.DATA, "text/plain",
                        ["log", "日志"], True, "日志文件"),

    # ==================== 字体 ====================
    "ttf": FileTypeInfo("ttf", FileCategory.FONT, "font/ttf",
                        ["ttf", "字体", "font"], False, "TrueType 字体"),
    "otf": FileTypeInfo("otf", FileCategory.FONT, "font/otf",
                        ["otf"], False, "OpenType 字体"),
    "woff": FileTypeInfo("woff", FileCategory.FONT, "font/woff",
                         ["woff"], False, "WOFF 字体"),
    "woff2": FileTypeInfo("woff2", FileCategory.FONT, "font/woff2",
                          ["woff2"], False, "WOFF2 字体"),

    # ==================== 设计/工程 ====================
    "psd": FileTypeInfo("psd", FileCategory.DESIGN, "image/vnd.adobe.photoshop",
                        ["psd", "photoshop"], False, "Photoshop 文件"),
    "ai": FileTypeInfo("ai", FileCategory.DESIGN, "application/illustrator",
                       ["ai", "illustrator"], False, "Illustrator 文件"),
    "sketch": FileTypeInfo("sketch", FileCategory.DESIGN, "application/sketch",
                           ["sketch"], False, "Sketch 文件"),
    "fig": FileTypeInfo("fig", FileCategory.DESIGN, "application/figma",
                        ["fig", "figma"], False, "Figma 文件"),
    "dwg": FileTypeInfo("dwg", FileCategory.DESIGN, "application/acad",
                        ["dwg", "cad", "autocad"], False, "AutoCAD 文件"),
    "obj": FileTypeInfo("obj", FileCategory.DESIGN, "model/obj",
                        ["obj", "3d", "模型"], False, "3D 模型"),
    "fbx": FileTypeInfo("fbx", FileCategory.DESIGN, "application/octet-stream",
                        ["fbx"], False, "FBX 3D 模型"),
    "blend": FileTypeInfo("blend", FileCategory.DESIGN, "application/x-blender",
                          ["blend", "blender"], False, "Blender 文件"),
}

# 输出目录 - 使用统一配置
from config import get_outputs_dir
OUTPUTS_DIR = get_outputs_dir()


# ============================================================
# 工具函数
# ============================================================

def generate_unique_filename(prefix: str, extension: str) -> str:
    """生成唯一文件名"""
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    short_id = str(uuid.uuid4())[:8]
    return f"{prefix}_{timestamp}_{short_id}.{extension}"


def _format_size(size_bytes: int) -> str:
    """格式化文件大小"""
    if size_bytes < 1024:
        return f"{size_bytes} B"
    elif size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.1f} KB"
    else:
        return f"{size_bytes / (1024 * 1024):.1f} MB"


def detect_file_type(skill_name: str, skill_description: str = "", params: Dict = None) -> FileTypeInfo:
    """
    根据技能名称、描述和参数检测目标文件类型

    Args:
        skill_name: 技能名称
        skill_description: 技能描述
        params: 执行参数（可能包含 output_type 指定）

    Returns:
        FileTypeInfo: 匹配到的文件类型信息
    """
    print(f"\n[FileType] Detecting file type for skill: {skill_name}")
    print(f"[FileType] Description: {skill_description[:100] if skill_description else 'None'}...")

    # 1. 优先使用显式指定的输出类型
    if params:
        explicit_type = params.get("output_type") or params.get("outputType") or params.get("file_type")
        if explicit_type:
            explicit_type = explicit_type.lower().strip(".")
            if explicit_type in FILE_TYPES:
                print(f"[FileType] Using explicit type from params: {explicit_type}")
                return FILE_TYPES[explicit_type]

    # 2. 构建搜索文本
    combined = f"{skill_name or ''} {skill_description or ''}".lower()
    print(f"[FileType] Combined search text: {combined[:100]}...")

    # 3. 精确匹配扩展名（优先级最高）
    for ext, info in FILE_TYPES.items():
        # 检查是否包含精确的扩展名引用
        if f".{ext}" in combined or f" {ext} " in combined or combined.endswith(f" {ext}") or combined.startswith(f"{ext} "):
            print(f"[FileType] Exact extension match: {ext}")
            return info

    # 3.5. 高优先级数据格式检测 - 这些格式应该优先于 HTML
    # 当用户明确提到这些格式时，不应该默认生成 HTML
    high_priority_formats = [
        ("json", ["json", "转换成json", "转成json", "导出json", "生成json"]),
        ("xlsx", ["excel", "xlsx", "xls", "电子表格"]),
        ("csv", ["csv"]),
        ("xml", ["xml"]),
        ("yaml", ["yaml", "yml"]),
        ("txt", ["txt", "文本文件", "纯文本"]),
        ("md", ["markdown", "md文件"]),
    ]
    for ext, keywords in high_priority_formats:
        for kw in keywords:
            if kw in combined:
                print(f"[FileType] High-priority format detected: {ext} (keyword: {kw})")
                return FILE_TYPES[ext]

    # 4. 关键词匹配（使用单词边界匹配，避免子字符串误匹配）
    # 例如: "understanding" 不应匹配 "rs"
    def is_word_match(keyword: str, text: str) -> bool:
        """检查关键词是否作为独立单词存在于文本中"""
        # 对于中文关键词，直接使用子字符串匹配
        if any('\u4e00' <= c <= '\u9fff' for c in keyword):
            return keyword in text
        # 对于英文关键词，使用单词边界匹配
        pattern = r'(?:^|[\s\-_.,;:!?()[\]{}\'"/])' + re.escape(keyword) + r'(?:$|[\s\-_.,;:!?()[\]{}\'"/])'
        return bool(re.search(pattern, text, re.IGNORECASE))

    best_match = None
    best_score = 0
    match_details = []

    for ext, info in FILE_TYPES.items():
        score = 0
        matched_keywords = []
        for keyword in info.keywords:
            if is_word_match(keyword, combined):
                # 关键词越长，匹配越精确，分数越高
                score += len(keyword)
                matched_keywords.append(keyword)

        if score > 0:
            match_details.append((ext, score, matched_keywords))

        if score > best_score:
            best_score = score
            best_match = info

    # 打印匹配详情
    if match_details:
        match_details.sort(key=lambda x: x[1], reverse=True)
        print(f"[FileType] Keyword matches (top 5):")
        for ext, score, keywords in match_details[:5]:
            print(f"[FileType]   - {ext}: score={score}, keywords={keywords}")

    if best_match:
        print(f"[FileType] Best match: {best_match.extension} (score={best_score})")
    else:
        print(f"[FileType] No match found, defaulting to HTML")

    # 5. 默认返回 HTML
    return best_match or FILE_TYPES["html"]


def get_supported_types() -> Dict[str, List[str]]:
    """获取所有支持的文件类型，按类别分组"""
    result = {}
    for ext, info in FILE_TYPES.items():
        category = info.category.value
        if category not in result:
            result[category] = []
        result[category].append({
            "extension": ext,
            "description": info.description,
            "can_generate": info.can_generate
        })
    return result


# 不需要输出文件的技能关键词
NO_OUTPUT_KEYWORDS = [
    # 对话/问答类
    "问答", "聊天", "对话", "chat", "qa", "assistant", "助手", "客服",
    "回答", "answer", "conversation", "bot", "机器人",
    # 查询/搜索类
    "查询", "搜索", "search", "query", "lookup", "find",
    # 通知/提醒类
    "通知", "提醒", "notify", "remind", "alert",
    # 验证/检查类
    "验证", "校验", "检查", "validate", "check", "verify",
    # 翻译类（通常直接显示结果）
    "翻译", "translate", "translation",
]

# 需要输出文件的技能关键词
OUTPUT_KEYWORDS = [
    # 生成/创建类
    "生成", "创建", "generate", "create", "制作", "build",
    # 分析/报表类
    "分析", "报表", "报告", "analyze", "analysis", "report", "统计",
    # 导出类
    "导出", "export", "下载", "download",
    # 文档类
    "文档", "document", "doc", "word", "pdf",
    # 表格类
    "表格", "excel", "spreadsheet", "csv",
    # 图片/设计类
    "图片", "图像", "image", "picture", "设计", "design", "画", "draw", "canvas",
    # 演示类
    "ppt", "演示", "slides", "presentation",
    # 代码类
    "代码", "code", "脚本", "script",
]


def should_generate_output(skill_name: str, skill_description: str = "") -> bool:
    """
    根据技能名称和描述判断是否应该生成输出文件

    Args:
        skill_name: 技能名称
        skill_description: 技能描述

    Returns:
        bool: True 表示应该生成输出文件，False 表示不需要
    """
    combined = f"{skill_name or ''} {skill_description or ''}".lower()

    # 计算匹配分数
    no_output_score = 0
    output_score = 0

    matched_no_output = []
    matched_output = []

    for keyword in NO_OUTPUT_KEYWORDS:
        if keyword.lower() in combined:
            no_output_score += len(keyword)
            matched_no_output.append(keyword)

    for keyword in OUTPUT_KEYWORDS:
        if keyword.lower() in combined:
            output_score += len(keyword)
            matched_output.append(keyword)

    print(f"[OutputDetect] Skill: {skill_name}")
    print(f"[OutputDetect] No-output keywords matched: {matched_no_output} (score: {no_output_score})")
    print(f"[OutputDetect] Output keywords matched: {matched_output} (score: {output_score})")

    # 如果有明确的输出关键词，优先生成文件
    if output_score > 0 and output_score >= no_output_score:
        print(f"[OutputDetect] Decision: GENERATE output (output keywords stronger)")
        return True

    # 如果只有无输出关键词匹配，则不生成
    if no_output_score > 0 and output_score == 0:
        print(f"[OutputDetect] Decision: NO output (no-output keywords matched)")
        return False

    # 默认生成输出文件
    print(f"[OutputDetect] Decision: GENERATE output (default)")
    return True


# ============================================================
# 文件生成器类
# ============================================================

class FileGenerator:
    """统一的文件生成器"""

    def __init__(self, output_dir: Path = None):
        self.output_dir = output_dir or OUTPUTS_DIR
        self.output_dir.mkdir(exist_ok=True)

    def generate(
        self,
        file_type: Union[str, FileTypeInfo],
        title: str,
        content: Any,
        extra_data: Dict = None
    ) -> Tuple[str, str, str]:
        """
        生成文件

        Args:
            file_type: 文件类型（扩展名或 FileTypeInfo）
            title: 标题
            content: 内容
            extra_data: 额外数据

        Returns:
            (filename, url, size_str)
        """
        if isinstance(file_type, str):
            file_type = FILE_TYPES.get(file_type.lower(), FILE_TYPES["txt"])

        print(f"[Generator] Generating {file_type.extension} for '{title}'")
        print(f"[Generator] Content type: {type(content)}")

        # 根据文件类型调用对应的生成方法
        method_name = f"_generate_{file_type.extension}"
        if hasattr(self, method_name):
            print(f"[Generator] Using method: {method_name}")
            return getattr(self, method_name)(title, content, extra_data)

        # 检查是否是可以作为文本处理的类型
        if file_type.category in [FileCategory.CODE, FileCategory.DATA]:
            print(f"[Generator] Using text file generator for category: {file_type.category}")
            return self._generate_text_file(title, content, file_type.extension, extra_data)

        # 不支持的类型，回退到 HTML
        print(f"[Generator] Falling back to HTML")
        return self._generate_html(title, str(content), extra_data)

    # ==================== 文档类生成器 ====================

    def _generate_txt(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成 TXT 文件"""
        filename = generate_unique_filename("output", "txt")
        filepath = self.output_dir / filename

        txt_content = f"""{title}
{"=" * len(title)}

{content}

---
生成时间: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}
"""
        filepath.write_text(txt_content, encoding="utf-8")
        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)

    def _generate_md(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成 Markdown 文件"""
        filename = generate_unique_filename("output", "md")
        filepath = self.output_dir / filename

        md_content = f"""# {title}

{content}

---
*生成时间: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}*
"""
        filepath.write_text(md_content, encoding="utf-8")
        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)

    def _generate_docx(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成 Word 文档"""
        if not HAS_DOCX:
            print(f"[Docx] python-docx not installed, falling back to HTML")
            return self._generate_html(title, content, extra_data)

        filename = generate_unique_filename("output", "docx")
        filepath = self.output_dir / filename

        print(f"[Docx] Generating Word document: {filename}")
        print(f"[Docx] Title: {title}")
        print(f"[Docx] Content type: {type(content)}, length: {len(str(content)) if content else 0}")

        try:
            doc = Document()
            doc.add_heading(title, level=0)

            # 处理内容
            if content:
                # 确保 content 是字符串
                if isinstance(content, (dict, list)):
                    content = json.dumps(content, ensure_ascii=False, indent=2)
                elif not isinstance(content, str):
                    content = str(content)

                paragraphs = content.split('\n\n')
                for para in paragraphs:
                    para = para.strip()
                    if not para:
                        continue

                    try:
                        if para.startswith('# '):
                            doc.add_heading(para[2:], level=1)
                        elif para.startswith('## '):
                            doc.add_heading(para[3:], level=2)
                        elif para.startswith('### '):
                            doc.add_heading(para[4:], level=3)
                        elif para.startswith('- ') or para.startswith('* '):
                            lines = para.split('\n')
                            for line in lines:
                                line = line.strip()
                                if line.startswith('- ') or line.startswith('* '):
                                    doc.add_paragraph(line[2:], style='List Bullet')
                                elif line:
                                    doc.add_paragraph(line)
                        elif para.startswith('```'):
                            code_content = para.strip('`').split('\n', 1)
                            if len(code_content) > 1:
                                code_para = doc.add_paragraph()
                                code_para.style = 'No Spacing'
                                run = code_para.add_run(code_content[1].rstrip('`'))
                                run.font.name = 'Consolas'
                                run.font.size = DocxPt(10)
                        else:
                            doc.add_paragraph(para)
                    except Exception as e:
                        print(f"[Docx] Error processing paragraph: {e}")
                        # 直接添加为普通段落
                        doc.add_paragraph(para)

            doc.add_paragraph()
            doc.add_paragraph(f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

            doc.save(str(filepath))

            # 验证文件
            size = filepath.stat().st_size
            if size < 1000:  # docx 文件应该至少有几 KB
                print(f"[Docx] WARNING: File seems too small ({size} bytes)")

            print(f"[Docx] Successfully generated: {filename}, size: {size}")
            return filename, f"/outputs/{filename}", _format_size(size)

        except Exception as e:
            print(f"[Docx] Error generating docx: {e}")
            import traceback
            traceback.print_exc()
            # 回退到 HTML
            return self._generate_html(title, str(content) if content else "", extra_data)

    def _generate_pdf(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成 PDF 文件"""
        if not HAS_REPORTLAB:
            return self._generate_html(title, content, extra_data)

        filename = generate_unique_filename("output", "pdf")
        filepath = self.output_dir / filename

        doc = SimpleDocTemplate(str(filepath), pagesize=A4)
        styles = getSampleStyleSheet()

        # 尝试注册中文字体
        try:
            font_paths = [
                "C:/Windows/Fonts/msyh.ttc",
                "C:/Windows/Fonts/simsun.ttc",
                "C:/Windows/Fonts/simhei.ttf",
                "/System/Library/Fonts/PingFang.ttc",  # macOS
                "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf",  # Linux
            ]
            font_registered = False
            for font_path in font_paths:
                if Path(font_path).exists():
                    pdfmetrics.registerFont(TTFont('Chinese', font_path))
                    font_registered = True
                    break

            if font_registered:
                styles.add(ParagraphStyle(name='ChineseNormal', fontName='Chinese', fontSize=12, leading=18))
                styles.add(ParagraphStyle(name='ChineseTitle', fontName='Chinese', fontSize=18, leading=24, spaceAfter=20))
                body_style = styles['ChineseNormal']
                title_style = styles['ChineseTitle']
            else:
                body_style = styles['Normal']
                title_style = styles['Title']
        except Exception:
            body_style = styles['Normal']
            title_style = styles['Title']

        story = []
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 0.3 * inch))

        paragraphs = content.split('\n\n') if content else [""]
        for para in paragraphs:
            if para.strip():
                para_html = para.replace('\n', '<br/>')
                story.append(Paragraph(para_html, body_style))
                story.append(Spacer(1, 0.15 * inch))

        story.append(Spacer(1, 0.5 * inch))
        story.append(Paragraph(f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", styles['Normal']))

        doc.build(story)
        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)

    def _generate_rtf(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成 RTF 文件"""
        filename = generate_unique_filename("output", "rtf")
        filepath = self.output_dir / filename

        # 简单的 RTF 格式
        rtf_content = r"{\rtf1\ansi\deff0"
        rtf_content += r"{\fonttbl{\f0 Arial;}}"
        rtf_content += r"\f0\fs24 "
        rtf_content += r"\b " + title.replace('\\', '\\\\').replace('{', '\\{').replace('}', '\\}') + r"\b0\par\par"
        rtf_content += content.replace('\\', '\\\\').replace('{', '\\{').replace('}', '\\}').replace('\n', r'\par ')
        rtf_content += r"\par\par ---\par Generated: " + datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        rtf_content += r"}"

        filepath.write_text(rtf_content, encoding="utf-8")
        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)

    # ==================== 表格类生成器 ====================

    def _generate_xlsx(self, title: str, data: Any, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成 Excel 文件"""
        filename = generate_unique_filename("output", "xlsx")
        filepath = self.output_dir / filename

        print(f"[Excel] Generating Excel file: {filename}")
        print(f"[Excel] Input data type: {type(data)}")

        df = None

        if isinstance(data, pd.DataFrame):
            df = data
            print(f"[Excel] Using DataFrame directly, shape: {df.shape}")
        elif isinstance(data, list) and len(data) > 0:
            if isinstance(data[0], dict):
                df = pd.DataFrame(data)
                print(f"[Excel] Created DataFrame from list of dicts, shape: {df.shape}")
            elif isinstance(data[0], list):
                if len(data) > 1:
                    df = pd.DataFrame(data[1:], columns=data[0])
                else:
                    df = pd.DataFrame(data)
                print(f"[Excel] Created DataFrame from list of lists, shape: {df.shape}")
            else:
                df = pd.DataFrame({"数据": data})
                print(f"[Excel] Created DataFrame from simple list, shape: {df.shape}")
        elif isinstance(data, dict):
            # 检查是否是嵌套的数据结构
            if any(isinstance(v, (list, dict)) for v in data.values()):
                # 尝试将字典转换为表格格式
                try:
                    df = pd.DataFrame([data])
                except Exception:
                    # 如果失败，转换为简单的键值对
                    df = pd.DataFrame([(k, str(v)) for k, v in data.items()], columns=["键", "值"])
            else:
                df = pd.DataFrame([data])
            print(f"[Excel] Created DataFrame from dict, shape: {df.shape}")
        elif isinstance(data, str):
            # 尝试解析字符串
            data_str = data.strip()
            parsed = False

            # 尝试解析为 JSON
            if data_str.startswith('[') or data_str.startswith('{'):
                try:
                    parsed_data = json.loads(data_str)
                    if isinstance(parsed_data, list):
                        df = pd.DataFrame(parsed_data)
                    elif isinstance(parsed_data, dict):
                        df = pd.DataFrame([parsed_data])
                    parsed = True
                    print(f"[Excel] Parsed string as JSON, shape: {df.shape}")
                except json.JSONDecodeError:
                    pass

            # 尝试解析为 CSV
            if not parsed and (',' in data_str or '\t' in data_str):
                try:
                    from io import StringIO
                    # 尝试逗号分隔
                    df = pd.read_csv(StringIO(data_str))
                    if df.shape[1] > 1:
                        parsed = True
                        print(f"[Excel] Parsed string as CSV, shape: {df.shape}")
                except Exception:
                    pass

            if not parsed:
                # 按行分割
                lines = data_str.split('\n')
                if len(lines) > 1:
                    df = pd.DataFrame({"内容": lines})
                else:
                    df = pd.DataFrame({"内容": [data_str]})
                print(f"[Excel] Created DataFrame from text lines, shape: {df.shape}")
        else:
            df = pd.DataFrame({"结果": [str(data)]})
            print(f"[Excel] Created DataFrame from unknown type, shape: {df.shape}")

        # 确保 DataFrame 不为空
        if df is None or df.empty:
            df = pd.DataFrame({"信息": ["无数据"]})
            print(f"[Excel] DataFrame was empty, using placeholder")

        # 清理数据：确保所有值都是可序列化的
        for col in df.columns:
            df[col] = df[col].apply(lambda x: str(x) if isinstance(x, (dict, list)) else x)

        # 处理列名：确保不包含非法字符
        df.columns = [str(col).replace('\n', ' ').replace('\r', ' ')[:255] for col in df.columns]

        sheet_name = extra_data.get("sheet_name", "Sheet1") if extra_data else "Sheet1"
        # 确保 sheet_name 有效
        sheet_name = str(sheet_name)[:31].replace('/', '-').replace('\\', '-').replace('*', '-').replace('?', '-').replace('[', '-').replace(']', '-')

        print(f"[Excel] DataFrame columns: {list(df.columns)}")
        print(f"[Excel] DataFrame shape: {df.shape}")
        print(f"[Excel] Sheet name: {sheet_name}")

        # 检查 openpyxl 是否可用
        if not HAS_OPENPYXL:
            print(f"[Excel] openpyxl not available, falling back to CSV")
            csv_filename = filename.replace('.xlsx', '.csv')
            csv_filepath = self.output_dir / csv_filename
            df.to_csv(csv_filepath, index=False, encoding='utf-8-sig')
            size = csv_filepath.stat().st_size
            return csv_filename, f"/outputs/{csv_filename}", _format_size(size)

        try:
            # 使用 openpyxl 直接创建工作簿（更可靠）
            from openpyxl import Workbook
            from openpyxl.utils.dataframe import dataframe_to_rows

            wb = Workbook()
            ws = wb.active
            ws.title = sheet_name

            # 写入数据
            for r_idx, row in enumerate(dataframe_to_rows(df, index=False, header=True), 1):
                for c_idx, value in enumerate(row, 1):
                    # 确保值可以写入 Excel
                    if value is None:
                        ws.cell(row=r_idx, column=c_idx, value="")
                    elif isinstance(value, (dict, list)):
                        ws.cell(row=r_idx, column=c_idx, value=str(value))
                    else:
                        try:
                            ws.cell(row=r_idx, column=c_idx, value=value)
                        except Exception:
                            ws.cell(row=r_idx, column=c_idx, value=str(value))

            # 保存文件
            wb.save(str(filepath))
            wb.close()

            print(f"[Excel] Successfully wrote Excel file with openpyxl Workbook")

            # 验证文件是否有效
            file_size = filepath.stat().st_size
            print(f"[Excel] File size: {file_size} bytes")

            if file_size < 2000:  # 有效的 xlsx 文件通常至少有几 KB
                print(f"[Excel] WARNING: File seems too small, verifying...")
                # 尝试重新打开验证
                try:
                    test_wb = openpyxl.load_workbook(str(filepath))
                    test_wb.close()
                    print(f"[Excel] File validation passed")
                except Exception as ve:
                    print(f"[Excel] File validation failed: {ve}")
                    raise ve

        except Exception as e:
            print(f"[Excel] Error writing with openpyxl: {e}")
            import traceback
            traceback.print_exc()

            # 尝试使用 pandas 的 to_excel
            try:
                print(f"[Excel] Trying pandas to_excel...")
                df.to_excel(filepath, index=False, engine='openpyxl', sheet_name=sheet_name)
                print(f"[Excel] Successfully wrote with pandas to_excel")
            except Exception as e2:
                print(f"[Excel] pandas to_excel also failed: {e2}")

                # 最后回退：保存为 CSV
                csv_filename = filename.replace('.xlsx', '.csv')
                csv_filepath = self.output_dir / csv_filename
                df.to_csv(csv_filepath, index=False, encoding='utf-8-sig')
                print(f"[Excel] Fallback to CSV: {csv_filename}")
                size = csv_filepath.stat().st_size
                return csv_filename, f"/outputs/{csv_filename}", _format_size(size)

        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)

    def _generate_csv(self, title: str, data: Any, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成 CSV 文件"""
        filename = generate_unique_filename("output", "csv")
        filepath = self.output_dir / filename

        if isinstance(data, pd.DataFrame):
            df = data
        elif isinstance(data, list) and len(data) > 0:
            if isinstance(data[0], dict):
                df = pd.DataFrame(data)
            elif isinstance(data[0], list):
                if len(data) > 1:
                    df = pd.DataFrame(data[1:], columns=data[0])
                else:
                    df = pd.DataFrame(data)
            else:
                df = pd.DataFrame({"数据": data})
        elif isinstance(data, dict):
            df = pd.DataFrame([data])
        else:
            df = pd.DataFrame({"内容": [str(data)]})

        df.to_csv(filepath, index=False, encoding="utf-8-sig")
        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)

    # ==================== 演示类生成器 ====================

    def _generate_pptx(self, title: str, content: Any, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成 PowerPoint 文件"""
        if not HAS_PPTX:
            return self._generate_html(title, str(content), extra_data)

        filename = generate_unique_filename("output", "pptx")
        filepath = self.output_dir / filename

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        title_slide_layout = prs.slide_layouts[0]
        content_slide_layout = prs.slide_layouts[1]

        # 标题页
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

        # 内容处理
        slides_data = None
        text_content = ""

        if isinstance(content, dict):
            if "slides" in content:
                slides_data = content["slides"]
            elif "content" in content:
                text_content = content["content"]
            else:
                text_content = json.dumps(content, ensure_ascii=False, indent=2)
        elif isinstance(content, list):
            if all(isinstance(item, dict) and "title" in item for item in content):
                slides_data = content
            else:
                text_content = str(content)
        else:
            text_content = str(content) if content else ""

        if slides_data:
            for slide_info in slides_data:
                slide = prs.slides.add_slide(content_slide_layout)
                slide.shapes.title.text = slide_info.get("title", "")
                slide.placeholders[1].text = slide_info.get("content", "")
        elif text_content:
            sections = text_content.split('\n\n')
            current_content = []
            current_title = "内容"

            for section in sections:
                section = section.strip()
                if not section:
                    continue

                lines = section.split('\n')
                first_line = lines[0].strip()

                if first_line.startswith('#'):
                    if current_content:
                        slide = prs.slides.add_slide(content_slide_layout)
                        slide.shapes.title.text = current_title
                        slide.placeholders[1].text = '\n'.join(current_content)
                        current_content = []

                    current_title = first_line.lstrip('#').strip()
                    if len(lines) > 1:
                        current_content = lines[1:]
                else:
                    current_content.append(section)
                    if len('\n'.join(current_content)) > 500:
                        slide = prs.slides.add_slide(content_slide_layout)
                        slide.shapes.title.text = current_title
                        slide.placeholders[1].text = '\n'.join(current_content)
                        current_content = []
                        current_title = "续"

            if current_content:
                slide = prs.slides.add_slide(content_slide_layout)
                slide.shapes.title.text = current_title
                slide.placeholders[1].text = '\n'.join(current_content)

        prs.save(str(filepath))
        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)

    # ==================== 图片类生成器 ====================

    def _generate_image(self, title: str, content: Any, extension: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        """通用图片生成器"""
        if not HAS_PIL:
            return self._generate_svg(title, content, extra_data)

        filename = generate_unique_filename("output", extension)
        filepath = self.output_dir / filename

        # 处理不同类型的输入
        if isinstance(content, bytes):
            # 二进制数据
            filepath.write_bytes(content)
        elif isinstance(content, str):
            if content.startswith("data:image"):
                # Base64 数据 URL
                header, data = content.split(",", 1)
                filepath.write_bytes(base64.b64decode(data))
            elif content.startswith("http"):
                # URL - 不支持下载，生成占位图
                return self._generate_placeholder_image(title, extension, extra_data)
            else:
                # 尝试作为 base64 解码
                try:
                    filepath.write_bytes(base64.b64decode(content))
                except Exception:
                    return self._generate_placeholder_image(title, extension, extra_data)
        elif isinstance(content, dict):
            # 检查是否有图片数据
            if "data" in content:
                return self._generate_image(title, content["data"], extension, extra_data)
            elif "base64" in content:
                return self._generate_image(title, content["base64"], extension, extra_data)
            else:
                return self._generate_placeholder_image(title, extension, extra_data)
        else:
            return self._generate_placeholder_image(title, extension, extra_data)

        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)

    def _generate_placeholder_image(self, title: str, extension: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成占位图片"""
        if HAS_PIL:
            filename = generate_unique_filename("output", extension)
            filepath = self.output_dir / filename

            # 创建一个简单的占位图
            width = extra_data.get("width", 800) if extra_data else 800
            height = extra_data.get("height", 600) if extra_data else 600

            img = PILImage.new('RGB', (width, height), color=(240, 240, 240))

            # 保存
            if extension.lower() in ['jpg', 'jpeg']:
                img.save(filepath, 'JPEG', quality=90)
            elif extension.lower() == 'png':
                img.save(filepath, 'PNG')
            elif extension.lower() == 'gif':
                img.save(filepath, 'GIF')
            elif extension.lower() == 'bmp':
                img.save(filepath, 'BMP')
            elif extension.lower() == 'webp':
                img.save(filepath, 'WEBP')
            else:
                img.save(filepath)

            size = filepath.stat().st_size
            return filename, f"/outputs/{filename}", _format_size(size)
        else:
            # 回退到 SVG
            return self._generate_svg(title, None, extra_data)

    def _generate_png(self, title: str, content: Any, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_image(title, content, "png", extra_data)

    def _generate_jpg(self, title: str, content: Any, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_image(title, content, "jpg", extra_data)

    def _generate_gif(self, title: str, content: Any, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_image(title, content, "gif", extra_data)

    def _generate_bmp(self, title: str, content: Any, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_image(title, content, "bmp", extra_data)

    def _generate_webp(self, title: str, content: Any, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_image(title, content, "webp", extra_data)

    def _generate_tiff(self, title: str, content: Any, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_image(title, content, "tiff", extra_data)

    def _generate_ico(self, title: str, content: Any, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_image(title, content, "ico", extra_data)

    def _generate_svg(self, title: str, content: Any, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成 SVG 文件"""
        filename = generate_unique_filename("output", "svg")
        filepath = self.output_dir / filename

        if isinstance(content, str) and content.strip().startswith("<svg"):
            # 已经是 SVG 内容
            svg_content = content
        else:
            # 生成简单的占位 SVG
            width = extra_data.get("width", 400) if extra_data else 400
            height = extra_data.get("height", 300) if extra_data else 300
            svg_content = f'''<?xml version="1.0" encoding="UTF-8"?>
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {width} {height}" width="{width}" height="{height}">
  <rect width="100%" height="100%" fill="#f0f0f0"/>
  <text x="50%" y="45%" text-anchor="middle" font-family="Arial, sans-serif" font-size="20" fill="#666">{title}</text>
  <text x="50%" y="55%" text-anchor="middle" font-family="Arial, sans-serif" font-size="12" fill="#999">SVG 图形</text>
</svg>'''

        filepath.write_text(svg_content, encoding="utf-8")
        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)

    # ==================== 压缩类生成器 ====================

    def _generate_zip(self, title: str, content: Any, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成 ZIP 压缩包"""
        filename = generate_unique_filename("output", "zip")
        filepath = self.output_dir / filename

        with zipfile.ZipFile(filepath, 'w', zipfile.ZIP_DEFLATED) as zf:
            if isinstance(content, dict):
                # 字典 - 每个键值对作为一个文件
                for name, data in content.items():
                    if isinstance(data, (str, bytes)):
                        zf.writestr(name, data if isinstance(data, bytes) else data.encode('utf-8'))
                    else:
                        zf.writestr(f"{name}.json", json.dumps(data, ensure_ascii=False, indent=2))
            elif isinstance(content, list):
                # 列表 - 每个元素作为一个文件
                for i, item in enumerate(content):
                    if isinstance(item, dict) and "name" in item and "content" in item:
                        zf.writestr(item["name"], item["content"] if isinstance(item["content"], bytes) else str(item["content"]).encode('utf-8'))
                    else:
                        zf.writestr(f"file_{i}.txt", str(item).encode('utf-8'))
            elif isinstance(content, str):
                # 字符串 - 作为单个文件
                zf.writestr(f"{title}.txt", content.encode('utf-8'))
            else:
                zf.writestr(f"{title}.txt", str(content).encode('utf-8'))

        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)

    def _generate_tar(self, title: str, content: Any, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成 TAR 归档"""
        filename = generate_unique_filename("output", "tar")
        filepath = self.output_dir / filename

        import io

        with tarfile.open(filepath, 'w') as tf:
            if isinstance(content, dict):
                for name, data in content.items():
                    data_bytes = data if isinstance(data, bytes) else str(data).encode('utf-8')
                    tarinfo = tarfile.TarInfo(name=name)
                    tarinfo.size = len(data_bytes)
                    tf.addfile(tarinfo, io.BytesIO(data_bytes))
            else:
                data_bytes = str(content).encode('utf-8')
                tarinfo = tarfile.TarInfo(name=f"{title}.txt")
                tarinfo.size = len(data_bytes)
                tf.addfile(tarinfo, io.BytesIO(data_bytes))

        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)

    def _generate_gz(self, title: str, content: Any, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成 GZIP 压缩文件"""
        filename = generate_unique_filename("output", "gz")
        filepath = self.output_dir / filename

        data = content if isinstance(content, bytes) else str(content).encode('utf-8')
        with gzip.open(filepath, 'wb') as f:
            f.write(data)

        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)

    # ==================== 代码/网页生成器 ====================

    def _strip_markdown_code_block(self, content: str) -> str:
        """
        剥离 markdown 代码块标记
        例如: ```html\n<html>...</html>\n``` -> <html>...</html>
        也处理没有闭合标记的情况（AI 输出被截断时）
        """
        if not content:
            return content

        stripped = content.strip()

        # 检查是否以 ``` 开头
        if stripped.startswith("```"):
            # 找到第一个换行符（跳过语言标识如 ```html）
            first_newline = stripped.find("\n")

            if first_newline != -1:
                # 尝试找到结尾的 ```
                last_fence = stripped.rfind("\n```")

                if last_fence > first_newline:
                    # 正常情况：有闭合的代码块标记
                    return stripped[first_newline + 1:last_fence].strip()
                else:
                    # AI 输出被截断，没有闭合标记，直接去掉开头的 ```xxx\n
                    return stripped[first_newline + 1:].strip()

        return content

    def _generate_html(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成 HTML 文件"""
        filename = generate_unique_filename("output", "html")
        filepath = self.output_dir / filename

        # 先清理 markdown 代码块包裹
        if content:
            content = self._strip_markdown_code_block(content)

        # 检查是否已经是完整的 HTML
        if content and content.strip().lower().startswith(("<!doctype", "<html")):
            html_content = content
        else:
            html_content = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            min-height: 100vh;
            padding: 40px;
            color: #333;
        }}
        .container {{
            max-width: 900px;
            margin: 0 auto;
            background: white;
            border-radius: 16px;
            padding: 40px;
            box-shadow: 0 20px 60px rgba(0,0,0,0.2);
        }}
        h1 {{
            color: #1e293b;
            margin-bottom: 24px;
            font-size: 28px;
            border-bottom: 2px solid #e2e8f0;
            padding-bottom: 16px;
        }}
        .content {{
            line-height: 1.8;
            color: #475569;
        }}
        .meta {{
            margin-top: 30px;
            padding-top: 20px;
            border-top: 1px solid #e2e8f0;
            font-size: 12px;
            color: #94a3b8;
        }}
        pre {{
            background: #f8fafc;
            padding: 16px;
            border-radius: 8px;
            overflow-x: auto;
            margin: 16px 0;
        }}
        code {{
            font-family: 'JetBrains Mono', 'Fira Code', Consolas, monospace;
            font-size: 13px;
        }}
    </style>
</head>
<body>
    <div class="container">
        <h1>{title}</h1>
        <div class="content">
            {content}
        </div>
        <div class="meta">
            生成时间: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")} | AI Skills Platform
        </div>
    </div>
</body>
</html>"""

        filepath.write_text(html_content, encoding="utf-8")
        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)

    def _generate_text_file(self, title: str, content: str, extension: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成通用文本文件（代码、数据等）"""
        filename = generate_unique_filename("output", extension)
        filepath = self.output_dir / filename

        # 处理内容
        if isinstance(content, (dict, list)):
            if extension in ['json']:
                text_content = json.dumps(content, ensure_ascii=False, indent=2)
            elif extension in ['yaml', 'yml'] and HAS_YAML:
                text_content = yaml.dump(content, allow_unicode=True, default_flow_style=False)
            else:
                text_content = str(content)
        else:
            text_content = str(content) if content else ""

        filepath.write_text(text_content, encoding="utf-8")
        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)

    def _generate_css(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_text_file(title, content, "css", extra_data)

    def _generate_js(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_text_file(title, content, "js", extra_data)

    def _generate_ts(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_text_file(title, content, "ts", extra_data)

    def _generate_py(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_text_file(title, content, "py", extra_data)

    def _generate_java(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_text_file(title, content, "java", extra_data)

    def _generate_go(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_text_file(title, content, "go", extra_data)

    def _generate_rs(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_text_file(title, content, "rs", extra_data)

    def _generate_c(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_text_file(title, content, "c", extra_data)

    def _generate_cpp(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_text_file(title, content, "cpp", extra_data)

    def _generate_php(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_text_file(title, content, "php", extra_data)

    def _generate_rb(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_text_file(title, content, "rb", extra_data)

    def _generate_sh(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_text_file(title, content, "sh", extra_data)

    def _generate_bat(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_text_file(title, content, "bat", extra_data)

    def _generate_ps1(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_text_file(title, content, "ps1", extra_data)

    def _generate_vue(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_text_file(title, content, "vue", extra_data)

    def _generate_jsx(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_text_file(title, content, "jsx", extra_data)

    def _generate_tsx(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        return self._generate_text_file(title, content, "tsx", extra_data)

    # ==================== 数据格式生成器 ====================

    def _generate_json(self, title: str, data: Any, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成 JSON 文件"""
        filename = generate_unique_filename("output", "json")
        filepath = self.output_dir / filename

        if isinstance(data, (dict, list)):
            output_data = data
        else:
            output_data = {
                "title": title,
                "data": data,
                "generated_at": datetime.now().isoformat()
            }

        filepath.write_text(json.dumps(output_data, ensure_ascii=False, indent=2, default=str), encoding="utf-8")
        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)

    def _generate_xml(self, title: str, data: Any, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成 XML 文件"""
        filename = generate_unique_filename("output", "xml")
        filepath = self.output_dir / filename

        def dict_to_xml(d, root_name="root"):
            def convert(data, parent_tag):
                xml_str = ""
                if isinstance(data, dict):
                    for key, value in data.items():
                        safe_key = str(key).replace(" ", "_")
                        xml_str += f"<{safe_key}>{convert(value, safe_key)}</{safe_key}>"
                elif isinstance(data, list):
                    for item in data:
                        xml_str += f"<item>{convert(item, 'item')}</item>"
                else:
                    xml_str = str(data).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                return xml_str

            return f'<?xml version="1.0" encoding="UTF-8"?>\n<{root_name}>{convert(d, root_name)}</{root_name}>'

        if isinstance(data, str) and data.strip().startswith("<?xml"):
            xml_content = data
        elif isinstance(data, (dict, list)):
            xml_content = dict_to_xml({"title": title, "data": data, "generated_at": datetime.now().isoformat()})
        else:
            xml_content = dict_to_xml({"title": title, "content": str(data)})

        filepath.write_text(xml_content, encoding="utf-8")
        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)

    def _generate_yaml(self, title: str, data: Any, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成 YAML 文件"""
        filename = generate_unique_filename("output", "yaml")
        filepath = self.output_dir / filename

        if HAS_YAML:
            if isinstance(data, (dict, list)):
                yaml_content = yaml.dump(data, allow_unicode=True, default_flow_style=False, sort_keys=False)
            else:
                yaml_content = yaml.dump({"title": title, "content": str(data)}, allow_unicode=True)
        else:
            # 简单的 YAML 格式化
            yaml_content = f"# {title}\n"
            if isinstance(data, dict):
                for k, v in data.items():
                    yaml_content += f"{k}: {v}\n"
            else:
                yaml_content += f"content: {data}\n"

        filepath.write_text(yaml_content, encoding="utf-8")
        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)

    def _generate_toml(self, title: str, data: Any, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成 TOML 文件"""
        filename = generate_unique_filename("output", "toml")
        filepath = self.output_dir / filename

        def to_toml(d, prefix=""):
            lines = []
            simple = {}
            complex_items = {}

            if isinstance(d, dict):
                for k, v in d.items():
                    if isinstance(v, dict):
                        complex_items[k] = v
                    elif isinstance(v, list):
                        if all(isinstance(i, (str, int, float, bool)) for i in v):
                            simple[k] = v
                        else:
                            complex_items[k] = v
                    else:
                        simple[k] = v

                for k, v in simple.items():
                    if isinstance(v, str):
                        lines.append(f'{k} = "{v}"')
                    elif isinstance(v, bool):
                        lines.append(f'{k} = {str(v).lower()}')
                    elif isinstance(v, list):
                        lines.append(f'{k} = {json.dumps(v)}')
                    else:
                        lines.append(f'{k} = {v}')

                for k, v in complex_items.items():
                    section = f"{prefix}.{k}" if prefix else k
                    lines.append(f"\n[{section}]")
                    lines.append(to_toml(v, section))
            else:
                lines.append(f'value = "{d}"')

            return "\n".join(lines)

        toml_content = f"# {title}\n" + to_toml(data if isinstance(data, dict) else {"content": data})
        filepath.write_text(toml_content, encoding="utf-8")
        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)

    def _generate_ini(self, title: str, data: Any, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成 INI 配置文件"""
        filename = generate_unique_filename("output", "ini")
        filepath = self.output_dir / filename

        lines = [f"; {title}", f"; Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", ""]

        if isinstance(data, dict):
            for section, values in data.items():
                lines.append(f"[{section}]")
                if isinstance(values, dict):
                    for k, v in values.items():
                        lines.append(f"{k} = {v}")
                else:
                    lines.append(f"value = {values}")
                lines.append("")
        else:
            lines.append("[default]")
            lines.append(f"content = {data}")

        filepath.write_text("\n".join(lines), encoding="utf-8")
        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)

    def _generate_sql(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成 SQL 文件"""
        filename = generate_unique_filename("output", "sql")
        filepath = self.output_dir / filename

        sql_content = f"-- {title}\n-- Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n{content}"
        filepath.write_text(sql_content, encoding="utf-8")
        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)

    def _generate_log(self, title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
        """生成日志文件"""
        filename = generate_unique_filename("output", "log")
        filepath = self.output_dir / filename

        log_content = f"[{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}] {title}\n{'=' * 50}\n{content}"
        filepath.write_text(log_content, encoding="utf-8")
        size = filepath.stat().st_size
        return filename, f"/outputs/{filename}", _format_size(size)


# ============================================================
# 全局实例和兼容性函数
# ============================================================

_generator = FileGenerator()


def _save_raw_html(html_content: str) -> Tuple[str, str, str]:
    """直接保存原始 HTML（不使用模板包装）"""
    filename = generate_unique_filename("output", "html")
    filepath = OUTPUTS_DIR / filename
    filepath.write_text(html_content, encoding="utf-8")
    size = filepath.stat().st_size
    return filename, f"/outputs/{filename}", _format_size(size)


# 兼容旧接口
def generate_html_file(title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
    return _generator._generate_html(title, content, extra_data)

def generate_markdown_file(title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
    return _generator._generate_md(title, content, extra_data)

def generate_json_file(title: str, data: Dict, extra_data: Dict = None) -> Tuple[str, str, str]:
    return _generator._generate_json(title, data, extra_data)

def generate_csv_file(title: str, headers: list, rows: list, extra_data: Dict = None) -> Tuple[str, str, str]:
    data = [dict(zip(headers, row)) for row in rows]
    return _generator._generate_csv(title, data, extra_data)

def generate_excel_file(title: str, data: Any, extra_data: Dict = None) -> Tuple[str, str, str]:
    return _generator._generate_xlsx(title, data, extra_data)

def generate_txt_file(title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
    return _generator._generate_txt(title, content, extra_data)

def generate_pdf_file(title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
    return _generator._generate_pdf(title, content, extra_data)

def generate_pptx_file(title: str, content: str, slides_data: List[Dict] = None, extra_data: Dict = None) -> Tuple[str, str, str]:
    return _generator._generate_pptx(title, content if slides_data is None else {"slides": slides_data}, extra_data)

def generate_docx_file(title: str, content: str, extra_data: Dict = None) -> Tuple[str, str, str]:
    return _generator._generate_docx(title, content, extra_data)


def generate_output_file(
    skill_name: str,
    skill_description: str,
    execution_result: Any,
    execution_output: str,
    params: Dict = None,
    output_config: Dict = None  # 新增：技能输出配置
) -> Optional[Dict[str, str]]:
    """
    根据技能信息生成输出文件

    Args:
        skill_name: 技能名称
        skill_description: 技能描述
        execution_result: 执行结果
        execution_output: 执行输出
        params: 执行参数
        output_config: 技能输出配置（来自 skill.output_config）
            - enabled: 是否生成文件（默认 True）
            - preferred_type: 优先文件类型
            - filename_template: 文件名模板

    Returns:
        {"name": "文件名", "type": "类型", "url": "下载URL", "size": "大小"}
    """
    # 1. 先检查技能级别的输出配置
    if output_config:
        print(f"[OutputFile] Using output_config: {output_config}")
        # 如果配置明确禁用输出文件
        if output_config.get("enabled") is False:
            print(f"[OutputFile] Output disabled by skill config")
            return None
    else:
        # 没有配置时，根据关键词自动判断
        if not should_generate_output(skill_name, skill_description):
            print(f"[OutputFile] Output disabled by keyword detection")
            return None

    # 2. 检查 execution_result 中的特殊标记（向后兼容）
    if isinstance(execution_result, dict):
        # 检查是否明确不需要输出文件
        if execution_result.get("_no_output_file"):
            print(f"[OutputFile] Skill requested no output file via _no_output_file marker")
            return None

        # 检查是否有 _output_file 标记
        if "_output_file" in execution_result:
            output_info = execution_result["_output_file"]
            if "path" in output_info:
                try:
                    file_size = Path(output_info["path"]).stat().st_size
                    size_str = _format_size(file_size)
                except Exception:
                    size_str = output_info.get("size", "未知")
            else:
                size_str = output_info.get("size", "未知")

            return {
                "name": output_info.get("name", "output"),
                "type": output_info.get("type", "file"),
                "url": output_info.get("url", ""),
                "size": size_str
            }

        # 检查是否有 _html 标记
        if "_html" in execution_result:
            raw_html = execution_result["_html"]
            filename, url, size = _save_raw_html(raw_html)
            return {"name": filename, "type": "html", "url": url, "size": size}

        # 检查是否有 _dataframe 或 _output_excel 标记
        if "_dataframe" in execution_result or "_output_excel" in execution_result:
            excel_data = execution_result.get("_dataframe") or execution_result.get("_output_excel")
            filename, url, size = _generator._generate_xlsx(skill_name, excel_data)
            return {"name": filename, "type": "excel", "url": url, "size": size}

    # 检查是否是完整的 HTML
    raw_html = None
    if isinstance(execution_result, str):
        stripped = execution_result.strip().lower()
        if stripped.startswith("<!doctype") or stripped.startswith("<html"):
            raw_html = execution_result
    if not raw_html and execution_output:
        stripped = execution_output.strip().lower()
        if stripped.startswith("<!doctype") or stripped.startswith("<html"):
            raw_html = execution_output

    if raw_html:
        filename, url, size = _save_raw_html(raw_html)
        return {"name": filename, "type": "html", "url": url, "size": size}

    # 检测文件类型
    # 优先使用 output_config 中指定的类型
    if output_config and output_config.get("preferred_type"):
        preferred = output_config["preferred_type"].lower().strip(".")
        if preferred in FILE_TYPES:
            file_type = FILE_TYPES[preferred]
            print(f"[OutputFile] Using preferred type from config: {file_type.extension}")
        else:
            print(f"[OutputFile] Preferred type '{preferred}' not found, falling back to detection")
            file_type = detect_file_type(skill_name, skill_description, params)
    else:
        file_type = detect_file_type(skill_name, skill_description, params)
    print(f"[OutputFile] Final file type: {file_type.extension}")

    # 准备内容
    content = None
    if isinstance(execution_result, dict):
        print(f"[OutputFile] execution_result is dict with keys: {list(execution_result.keys())}")
        if "content" in execution_result:
            content = execution_result["content"]
            print(f"[OutputFile] Using 'content' field, type: {type(content)}")
        elif "data" in execution_result:
            content = execution_result["data"]
            print(f"[OutputFile] Using 'data' field, type: {type(content)}")
        elif "image" in execution_result or "base64" in execution_result:
            # 图片数据
            content = execution_result.get("image") or execution_result.get("base64")
            print(f"[OutputFile] Using image/base64 field")
        else:
            # 过滤特殊字段
            content = {k: v for k, v in execution_result.items() if not k.startswith("_")}
            print(f"[OutputFile] Using filtered dict as content, keys: {list(content.keys()) if isinstance(content, dict) else 'N/A'}")
    elif execution_result:
        content = execution_result
        print(f"[OutputFile] Using execution_result directly, type: {type(content)}")

    if not content and execution_output:
        content = execution_output
        print(f"[OutputFile] Using execution_output as content")

    if not content:
        content = "技能执行完成"
        print(f"[OutputFile] No content, using default message")

    print(f"[OutputFile] Final content type: {type(content)}")
    if isinstance(content, str):
        print(f"[OutputFile] Content preview: {content[:200]}...")
    elif isinstance(content, dict):
        print(f"[OutputFile] Content dict keys: {list(content.keys())}")

    # 生成文件
    try:
        print(f"[OutputFile] Generating file with type: {file_type.extension}")
        filename, url, size = _generator.generate(file_type, skill_name, content)
        print(f"[OutputFile] Generated: {filename}")
        return {
            "name": filename,
            "type": file_type.extension,
            "url": url,
            "size": size
        }
    except Exception as e:
        import traceback
        print(f"[OutputFile] ERROR generating file: {e}")
        traceback.print_exc()
        # 出错时回退到 HTML
        result_content = f"<p>执行结果：</p><pre>{content}</pre>" if content else "<p>执行完成</p>"
        filename, url, size = _generator._generate_html(f"{skill_name} - 结果", result_content)
        return {"name": filename, "type": "html", "url": url, "size": size}
